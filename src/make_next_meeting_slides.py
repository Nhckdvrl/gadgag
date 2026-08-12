#!/usr/bin/env python3
"""Generate the 2026-08-19 advisor-meeting deck from frozen project evidence."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables/meeting_20260819_lexical_arbitration.pptx"
FONT = "Noto Sans CJK JP"
NAVY = RGBColor(22, 46, 81)
BLUE = RGBColor(42, 111, 151)
TEAL = RGBColor(33, 139, 130)
RED = RGBColor(180, 61, 61)
GOLD = RGBColor(213, 157, 55)
LIGHT = RGBColor(240, 244, 248)
MID = RGBColor(100, 112, 126)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(25, 28, 32)


def rect(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def text(slide, value, x, y, w, h, size=20, color=BLACK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = value
    run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=18, color=BLACK, level_gap=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple): value, level = item
        else: value, level = item, 0
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = value; p.level = level; p.font.name = FONT
        p.font.size = Pt(size - level * 2); p.font.color.rgb = color
        p.space_after = Pt(level_gap); p.line_spacing = 1.05
    return box


def base(prs, title, number, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 10, .16, NAVY)
    text(slide, title, .45, .28, 8.9, .6, 24, NAVY, True)
    if subtitle: text(slide, subtitle, .47, .86, 8.9, .35, 11, MID)
    text(slide, str(number), 9.35, 7.05, .3, .2, 9, MID, align=PP_ALIGN.RIGHT)
    return slide


def footer(slide, source):
    text(slide, source, .45, 6.92, 8.6, .24, 8, MID)


def add_table(slide, rows, widths, x, y, h, header=True, font=12):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y),
                                   Inches(sum(widths)), Inches(h)).table
    for i, width in enumerate(widths): table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r,c); cell.text = str(value)
            cell.margin_left = Inches(.06); cell.margin_right = Inches(.04)
            cell.margin_top = Inches(.035); cell.margin_bottom = Inches(.025)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if header and r==0 else (LIGHT if r%2 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.name=FONT; p.font.size=Pt(font); p.font.color.rgb=WHITE if header and r==0 else BLACK
                if header and r==0: p.font.bold=True
    return table


def main() -> None:
    prs = Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide,0,0,10,7.5,WHITE)
    rect(slide,0,0,10,.22,NAVY); rect(slide,.55,1.15,.12,4.8,TEAL)
    text(slide,"文脈と言語慣習が衝突するとき",.95,1.28,8.2,.65,30,NAVY,True)
    text(slide,"多言語 LLM における語彙的仲裁の分解",.96,2.05,8.1,.55,23,BLUE,True)
    text(slide,"When Context and Language Disagree",.98,2.82,7.6,.38,17,MID)
    text(slide,"今回：scientific question の明確化・先行研究境界・五種類の controls",.98,4.15,8.1,.7,18,BLACK)
    text(slide,"M1  向　｜　2026/08/19",.98,6.42,4.5,.35,13,MID)

    s=base(prs,"前回の観察から scientific question へ",2)
    rect(s,.55,1.25,4.15,4.85,LIGHT,radius=True)
    text(s,"前回",.85,1.55,1.1,.35,19,MID,True)
    text(s,"誤答は二種類の failure を\n混ぜていないか？",.85,2.12,3.4,1.0,24,NAVY,True)
    bullets(s,["context を利用しない","正しい方向へ動くが最終的に誤る"],.95,3.35,3.35,1.35,17)
    rect(s,4.95,1.25,4.5,4.85,WHITE,TEAL,True)
    text(s,"今回：So what?",5.28,1.55,2.2,.35,19,TEAL,True)
    text(s,"cross-lingual lexical error は\nどの計算段階で生じるのか？",5.28,2.1,3.72,1.0,24,NAVY,True)
    bullets(s,["semantic-evidence extraction","lexical decision / arbitration"],5.35,3.35,3.6,1.35,18)
    text(s,"「FF が難しい」ではなく failure localization",5.35,5.22,3.55,.45,15,RED,True)

    s=base(prs,"先行研究：既に問われたことを contribution にしない",3)
    rows=[["Work","既に行ったこと","claim 不可","残る差分"],
          ["Stingray","FF/cognate・4 pairs・bias","FF が難しい／多言語化","Language×Sense の分離"],
          ["Tanwar","3 word types・context・incongruent","矛盾文／3 分類","stage-specific estimand"],
          ["Doppel","JC 3 tasks・shortcut・POS","JC error analysis","誤答内の context evidence"],
          ["RoDEval","WSD knowledge/bias/reliability","accuracy 不十分","exact-form conflict 分解"],
          ["Dumas / Oh","patching・context/lexical competition","表現分離／競争の初研究","cross-lingual collision"],]
    add_table(s,rows,[1.2,2.55,2.35,2.75],.55,1.25,4.95,font=10.5)
    footer(s,"Cahyawijaya+ 2025; Tanwar+ arXiv:2501.09127; Kitamura+ 2025; Zhang+ 2025; Dumas+ 2025; Oh+ 2026")

    s=base(prs,"Main RQ：抽出失敗か，仲裁失敗か",4)
    rect(s,.6,1.25,8.8,1.4,LIGHT,radius=True)
    text(s,"文脈が支持する語義と言語慣習が支持する語義が衝突するとき，\n誤りは文脈語義の抽出段階か，抽出後の最終語義選択か？",.95,1.56,8.05,.82,23,NAVY,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    for x,label,desc,color in [(0.75,"RQ1","Evidence extraction\nLanguage 固定・Sense のみ変更",BLUE),
                               (3.62,"RQ2","Collision specificity\n五種類の matched controls",TEAL),
                               (6.5,"RQ3","Arbitration mechanism\ntarget residual / MLP",GOLD)]:
        rect(s,x,3.15,2.55,2.25,WHITE,color,True)
        text(s,label,x+.2,3.4,.55,.35,18,color,True)
        text(s,desc,x+.2,3.95,2.15,.95,17,NAVY,True,PP_ALIGN.CENTER)
    text(s,"B = behavioral backbone　→　A = conditional mechanistic explanation",1.4,5.9,7.2,.45,18,RED,True,PP_ALIGN.CENTER)

    s=base(prs,"仮説：正しい evidence があっても答えは誤り得る",5)
    labels=[("Local semantic context",BLUE),("Target representation",TEAL),("Lexical arbitration",GOLD),("Final choice",NAVY)]
    xs=[.45,2.85,5.25,7.65]
    for (label,color),x in zip(labels,xs):
        rect(s,x,2.25,1.9,1.05,WHITE,color,True); text(s,label,x+.1,2.5,1.7,.5,16,color,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    for x in (2.35,4.75,7.15): text(s,"→",x,2.45,.45,.4,26,MID,True,PP_ALIGN.CENTER)
    rect(s,4.98,1.18,2.42,.62,LIGHT,RED,True)
    text(s,"Language convention / shared form",5.1,1.34,2.18,.28,12,RED,True,PP_ALIGN.CENTER)
    text(s,"↓",6.0,1.77,.35,.35,22,RED,True,PP_ALIGN.CENTER)
    text(s,"pilot が示唆する paradox",.65,4.15,2.5,.35,18,TEAL,True)
    text(s,"context evidence は正方向\nしかし decision boundary を越えない",.65,4.62,3.35,1.0,24,NAVY,True)
    text(s,"研究価値",5.15,4.15,1.4,.35,18,GOLD,True)
    bullets(s,["failure diagnosis が変わる","一般 WSD と collision を分ける","remedy の対象段階が変わる"],5.15,4.62,3.85,1.35,17)

    s=base(prs,"RQ1：Language × Sense の完全 2×2",6)
    rows=[["","Sense 1 context","Sense 2 context"],
          ["Language 1","L1 × S1","L1 × S2"],
          ["Language 2","L2 × S1","L2 × S2"]]
    add_table(s,rows,[1.55,2.15,2.15],.7,1.42,2.2,font=17)
    text(s,"Semantic effect",6.25,1.38,2.25,.35,18,BLUE,True)
    text(s,"同じ言語のまま\nS1 → S2",6.25,1.9,2.45,.75,23,NAVY,True)
    text(s,"Language effect",6.25,3.12,2.25,.35,18,RED,True)
    text(s,"同じ sense のまま\nL1 → L2",6.25,3.65,2.45,.75,23,NAVY,True)
    bullets(s,["exact same form","target mask","language-only","marker-matched unrelated","repeated gloss / natural data"],.82,4.25,4.9,1.75,16)
    footer(s,"Stingray の既存四格を再利用。ただし estimand と controls を事前固定。")

    s=base(prs,"現在の behavioral evidence：誤答内部に context signal",7)
    rows=[["Condition","CI > 0"],["Full context","92 / 96"],["Target masked","91 / 96"],
          ["Language only","4 positive / 2 negative"],["Matched unrelated","1 positive / 3 negative"]]
    add_table(s,rows,[3.1,1.55],.55,1.3,3.85,font=15)
    rect(s,5.55,1.3,3.85,3.85,LIGHT,radius=True)
    text(s,"Generality gates",5.9,1.62,2.8,.35,20,TEAL,True)
    bullets(s,["official chat: full 48/48","natural masked−unrelated: 82/96","ID–MS / ID–TL: 48/48","Doppel natural: 4/4 models"],5.88,2.18,3.0,2.2,17)
    rect(s,.7,5.35,8.6,.82,WHITE,RED,True)
    text(s,"Correct semantic evidence can be present without yielding\nthe correct lexical decision.",.92,5.49,8.18,.52,16,RED,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    footer(s,"96 cells は独立反復ではなく protocol-robustness cells。主解析は item-level hierarchical / bootstrap。")

    s=base(prs,"RQ2：五種類の controls（novelty ではなく identification）",8)
    rows=[["Type","Form","Semantics","役割"],
          ["False friend","same exact","different","target conflict"],
          ["True friend","same exact","same","shared-form baseline"],
          ["Language-specific","one language only","—","ordinary lexical processing"],
          ["Different-form translation","different","same","no form collision"],
          ["Monolingual polysemy","same / one language","different senses","ordinary WSD"]]
    add_table(s,rows,[2.2,1.7,1.75,3.05],.48,1.22,4.7,font=13)
    text(s,"先生の minimum baseline：最初の三種類",.65,6.05,4.3,.35,16,TEAL,True)
    text(s,"追加二種類：form sharing と ambiguity を分離",5.05,6.05,4.2,.35,16,BLUE,True)
    footer(s,"Oi & Saito 2009; Hsieh+ 2017; Tarin+ 2025/26 も language-specific controls を使用。")

    s=base(prs,"Outcome-blind matching：旧 control の失敗を設計で修正",9)
    rect(s,.55,1.25,3.95,4.75,LIGHT,radius=True)
    text(s,"旧 pool：FAIL",.85,1.55,2.2,.38,22,RED,True)
    text(s,"exact bilingual POS\n+ every covariate ≤ 1 SD",.85,2.12,3.1,.78,18,NAVY)
    text(s,"0 pairs",.85,3.22,2.5,.6,34,RED,True)
    text(s,"→ causal claim に使わない",.85,4.05,2.9,.35,17,BLACK)
    rect(s,4.8,1.25,4.65,4.75,WHITE,TEAL,True)
    text(s,"新 target-first design",5.12,1.55,3.1,.38,22,TEAL,True)
    bullets(s,["11,000 natural-context candidates","24 / 27 false friends retained","24 true + 24 translation","max |SMD| = .087 / .097","4× validation reservoir: 96/group","Qwen3/Gemma outcomes 未読"],5.1,2.08,3.75,2.95,17)
    text(s,"POS・frequency・ratio・length・tokenization・Qwen2.5 difficulty",5.12,5.25,3.75,.45,12,MID)
    footer(s,"Joint cardinality matching。辞書 overlap は gold ではないため bilingual validation 前は candidate。")

    s=base(prs,"追加した language-specific control",10)
    text(s,"操作的定義",.65,1.25,1.6,.35,19,TEAL,True)
    text(s,"NFKC exact form が相手言語辞書に存在しない二字語",.65,1.72,7.9,.45,22,NAVY,True)
    bullets(s,["自然文 pool：ZH 7,943 / JA 7,923","final：ZH 23 / JA 24，max |SMD|=.094/.097","4× reservoir：ZH 92 / JA 96，max |SMD|<.10","frequency / length / Qwen3・Gemma token / gloss / Qwen2.5 difficulty","辞書非掲載 ≠ 絶対不存在 → human validation 必須"],.82,2.5,7.9,2.55,18)
    rect(s,.75,5.45,8.35,.63,LIGHT,RED,True)
    text(s,"役割：cross-language LCE ではなく，通常の one-language context extraction / decision baseline",.95,5.64,7.95,.27,16,RED,True,PP_ALIGN.CENTER)
    footer(s,"この control を false-friend language effect と同じ estimand に無理に入れない。")

    s=base(prs,"RQ3：B が paradox を発見し，A が条件付きで説明する",11)
    if (ROOT/"figures/target_component_patching.png").exists():
        s.shapes.add_picture(str(ROOT/"figures/target_component_patching.png"), Inches(.45), Inches(1.23), width=Inches(5.6))
    rect(s,6.18,1.25,3.3,4.85,LIGHT,radius=True)
    text(s,"現在の解釈",6.48,1.58,2.45,.35,19,TEAL,True)
    bullets(s,["semantic profile は ordinary WSD-like","target residual で両 signal が causal","MLP は主要 language signal を再現","attention-only は不安定"],6.45,2.12,2.55,2.1,16)
    text(s,"claim しない",6.48,4.72,1.7,.35,17,RED,True)
    text(s,"homograph circuit / neuron\n人間脳との同一性",6.48,5.12,2.45,.62,15,BLACK)
    footer(s,"Matched human-valid controls を通過した場合のみ collision-specific excess を報告。")

    s=base(prs,"次週までの hard gates と判定規則",12)
    steps=[("1","Doppel 708 rows × 2 bilinguals"),("2","true/translation 192 × 2"),
           ("3","language-specific 188 × 2"),("4","human-valid set で再 matching"),
           ("5","≥20 FF, all |SMD|≤.10 のみ unlock")]
    for i,(n,label) in enumerate(steps):
        y=1.15+i*.9; rect(s,.75,y,.58,.58,TEAL,TEAL,True)
        text(s,n,.75,y+.08,.58,.32,18,WHITE,True,PP_ALIGN.CENTER)
        text(s,label,1.55,y+.08,6.85,.38,18,NAVY,True)
    rect(s,.72,5.82,8.45,.72,WHITE,RED,True)
    text(s,"失敗時：collision-specific mechanism を削除。\nB の behavioral decomposition は独立に評価。",.92,5.94,8.05,.45,14,RED,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    footer(s,"Target outcomes を見て matching・除外基準を変更しない。SAE/neuron/mitigation は gate 後。")

    s=base(prs,"今回，先生に確認したい点",13)
    bullets(s,["Main RQ を「context extraction vs lexical arbitration の failure localization」としてよいか",
               "B を主 contribution，matched causal A を条件付き第二 contribution としてよいか",
               "language-specific を one-language extraction/decision baseline とする設計でよいか",
               "中日 bilingual annotator 二名をどのように確保するか"],.85,1.38,8.05,3.5,21,level_gap=14)
    rect(s,.78,5.45,8.45,.72,LIGHT,TEAL,True)
    text(s,"示したいこと：モデルは文脈を理解していても，語彙的仲裁で負けて誤答し得る。",1.05,5.66,7.9,.3,20,TEAL,True,PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True,exist_ok=True); prs.save(OUT)
    print(OUT)


if __name__ == "__main__": main()
